#!/usr/bin/env python3
"""Build the final presentation on the BITS Capstone 15-minute template.

The template's ten slides and their titles are preserved exactly; only the body text,
the images and the speaker notes are ours. Timings in the notes add up to 15 minutes.

Usage:
    python3 scripts/docs/build_deck.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[2]
TEMPLATE = ROOT / "Capstone_Project_15_Min_Demo_Template (1).pptx"
FIGURES = ROOT / "docs" / "figures"
OUT = ROOT / "deliverables" / "ppt" / "09-TalkToJesus-Capstone-Final.pptx"

BODY_TOP = Inches(1.55)          # the template's body starts at 2.00in; tightened
BODY_HEIGHT = Inches(5.25)


def set_body(slide, lines, size=17, image=None, image_width=None):
    """Replace the body placeholder's text, optionally alongside a figure."""
    body = slide.placeholders[1]
    # All four of left/top/width/height must be set together. Setting only top and
    # height makes python-pptx write an explicit extent whose width defaults to 0,
    # which renders the placeholder invisible.
    body.left = Inches(0.44)
    body.top = BODY_TOP
    body.height = BODY_HEIGHT
    body.width = Inches(5.35) if image is not None else Inches(9.12)

    frame = body.text_frame
    frame.word_wrap = True
    frame.clear()

    for i, entry in enumerate(lines):
        text, level = entry if isinstance(entry, tuple) else (entry, 0)
        para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        para.text = text
        para.level = level
        para.space_after = Pt(7)
        for run in para.runs:
            run.font.size = Pt(size - 2 * level)
            if text.endswith(":") or text.isupper():
                run.font.bold = True

    if image is not None:
        path = FIGURES / image
        slide.shapes.add_picture(
            str(path),
            left=Inches(6.05),
            top=Inches(1.60),
            width=image_width or Inches(3.55),
        )


def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text.strip()


def main() -> None:
    prs = Presentation(TEMPLATE)
    s = prs.slides

    # ---- Slide 1: title -------------------------------------------------------
    boxes = [sh for sh in s[0].shapes if sh.has_text_frame]
    replacements = {
        "Capstone Project Title": "TalkToJesus",
        "Presented by:": "Presented by:\nManish Rachakonda (2023EBCS668)",
        "Under the guidance of": "Under the guidance of\nSwapnil Saurav",
        "BSc CS (Academic Year)": "BSc Computer Science (2025-2026)",
    }
    for box in boxes:
        current = box.text_frame.text
        for key, value in replacements.items():
            if current.startswith(key):
                box.text_frame.text = value
                for para in box.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(15)
                break
    # Give the project a subtitle line under the title.
    title = next(b for b in boxes if b.text_frame.text == "TalkToJesus")
    for para in title.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(30)
            run.font.bold = True
    sub = s[0].shapes.add_textbox(Inches(2.60), Inches(4.62), Inches(6.10), Inches(0.60))
    sub.text_frame.word_wrap = True
    sub.text_frame.text = "A Bilingual Voice-First AI Spiritual Companion"
    for run in sub.text_frame.paragraphs[0].runs:
        run.font.size = Pt(14)
        run.font.italic = True
    set_notes(s[0], """
[0:00-0:20]  Name, roll number, supervisor, one sentence on what it is:
"A user speaks a question in English or Telugu and hears a scripture-grounded reply
in a pastoral voice." Then move on - do not linger on the title slide.
""")

    # ---- Slide 2: problem -----------------------------------------------------
    set_body(s[1], [
        "Telugu Bible apps have millions of downloads - and zero conversation.",
        ("One-way communication - they present content; none answer a question", 1),
        ("Access is gated on a person - a pastor must be available, scheduled, approachable", 1),
        ("Guidance is not situation-aware - a reading plan does not know what you are facing", 1),
        ("Language and culture - general AI assistants are not scripture-grounded, and their", 1),
        ("Telugu is not idiomatic - addressing a believer as \"naa bidda\" (నా బిడ్డ) is not", 1),
        ("a translation problem", 1),
        "",
        "The gap: a believer with a question at 2 a.m., in Telugu, has nowhere to ask it.",
    ], size=16)
    set_notes(s[1], """
[0:20-1:30]  Lead with the one-liner - millions of downloads, zero conversation.
Then the four gaps, one sentence each. Land on the 2 a.m. line: that is the whole
motivation. Do not read the bullets verbatim.
""")

    # ---- Slide 3: objectives & scope ------------------------------------------
    set_body(s[2], [
        "Objectives (PO1-PO7):",
        ("PO1-PO2  Cross-platform app + secure backend with auth, conversation, subscriptions", 1),
        ("PO3-PO4  LLM for scripture-grounded replies; voice in and voice out", 1),
        ("PO5      Bilingual English/Telugu with culturally correct address", 1),
        ("PO6-PO7  Subscription monetisation; containerised deployment on GCP", 1),
        "Scope:",
        ("In - conversation, Bible reader, hymn library, subscriptions, admin console", 1),
        ("Out - app store publication, video avatar, community features, RAG over scripture", 1),
    ], size=16)
    set_notes(s[2], """
[1:30-2:20]  Do NOT read seven objectives. Group them as shown: app + backend,
AI + voice, bilingual, money + deployment. Spend the saved time on the "Out" line -
being explicit that there is no RAG is a credibility point, and an examiner may ask.
""")

    # ---- Slide 4: literature review -------------------------------------------
    set_body(s[3], [
        "Conversational?  Scripture-grounded?  Telugu?  Voice-first?",
        ("YouVersion Bible App - strong content, no dialogue", 1),
        ("Telugu Bible apps - Telugu and scripture, but static content only", 1),
        ("Pray.com - guided prayer and audio, no conversation, no Telugu", 1),
        ("General AI assistants - conversational, but not scripture-grounded, weak Telugu", 1),
        ("Glorify / Abide - devotional, no conversational AI, no Telugu", 1),
        "",
        "Each covers some. None covers all four. That intersection is the contribution.",
    ], size=16)
    set_notes(s[3], """
[2:20-3:10]  The point is the empty column, not the competitor list. Say the four
criteria once, then: "every one of these covers some of the four; none covers all
four at once." That sentence is the UVP - deliver it deliberately.
""")

    # ---- Slide 5: architecture ------------------------------------------------
    set_body(s[4], [
        "Three tiers, one stateless API:",
        ("Presentation - Flutter (iOS/Android, Riverpod) + the admin console", 1),
        ("Application - Express 5 / TypeScript, 28 endpoints, Cloud Run, scales from zero", 1),
        ("Data - Supabase PostgreSQL (8 tables) + OpenAI, ElevenLabs, Razorpay", 1),
        "",
        "One decision drives the rest: the API holds a service-tier key that",
        "bypasses row-level security - so authorization lives in application code.",
    ], size=15, image="print/fig-26-architecture.png", image_width=Inches(3.65))
    set_notes(s[5 - 1], """
[3:10-4:30]  The most technical slide - slow down. Walk the diagram top to bottom.
Then land the service-key point: RLS is on but unpolicied and bypassed, so every
user-scoped query must filter in code. That is an honest architectural weakness and
naming it before the examiner does is worth more than hiding it.
""")

    # ---- Slide 6: tools & technologies ----------------------------------------
    set_body(s[5], [
        "Language / Framework:  Dart + Flutter 3.38.6  |  TypeScript 5.9 + Express 5 on Node 20",
        "State / Validation:    Riverpod 2.5  |  Zod 4.1",
        "Database:              Supabase (PostgreSQL), 8 tables, no ORM",
        "AI services:           OpenAI whisper-1 (STT) + gpt-4o (LLM)  |  ElevenLabs (TTS)",
        "Payments / Identity:   Razorpay subscriptions  |  Google OAuth 2.0 -> HS256 JWT",
        "Infrastructure:        Docker, Cloud Build, Cloud Run (us-central1), Artifact Registry",
        "Observability:         Winston, Sentry, PostHog, per-stage latency instrumentation",
        "Quality:               Jest (66 tests) + flutter_test (75 tests) = 141, gated in CI",
    ], size=14)
    set_notes(s[5], """
[4:30-5:00]  Read the categories, not every item. Finish on the last line - 141 tests,
all passing, gated in CI. Then >>> CHECK THE CLOCK. You need seven clear minutes for
the demo. If you are past 5:15, skip straight to slide 7.
""")

    # ---- Slide 7: implementation / demo ---------------------------------------
    set_body(s[6], [
        "LIVE DEMO",
        ("1. Sign in with Google", 1),
        ("2. Home screen - switch EN -> TE (తెలుగు), every label changes, no restart", 1),
        ("3. Bible reader in Telugu - Genesis 1, then change translation", 1),
        ("4. Hymn library - 12 bundled hymns, plays with no network", 1),
        ("5. Speak a question - record, send, hear the reply", 1),
        ("6. Free tier exhausted -> paywall -> Razorpay UPI -> badge turns active", 1),
        ("7. Conversation history - the full turn with scripture citations", 1),
        ("8. Admin console - live metrics, then the latency breakdown", 1),
        "",
        "Fallback: the recorded demo, and the console's offline demo mode.",
    ], size=14, image="print/fig-03-home-english.jpg", image_width=Inches(1.62))
    set_notes(s[6], """
[5:00-12:00]  SEVEN MINUTES. This is the slide that decides the viva.

Order matters: bilingual switch first (it is instant and visual), then Bible, then
hymns, then the voice turn. START THE VOICE TURN EARLY - it takes ~28 seconds to
come back. Fill that silence deliberately: "while that runs, note that every stage
is being timed - I will show you the numbers in two minutes." Do not apologise for
the wait; frame it as the thing you measured.

Then paywall -> UPI -> badge. Then history. Then the console.

FALLBACKS, in order: (1) if the network is slow, play the recorded demo from
10-Demo-Video-Links.txt; (2) if the backend is down, the console has an offline
demo mode - but SAY it is demo data; (3) if all else fails, walk the figures in
the report.
""")

    # ---- Slide 8: results -----------------------------------------------------
    set_body(s[7], [
        "20/20 functional requirements implemented  |  141 automated tests, 0 failing",
        "Full subscription flow verified end to end on device, real UPI payment",
        "",
        "MEASURED pipeline latency (live instance, p50):",
        ("Speech to text (Whisper)      4,004 ms    15%", 1),
        ("AI response (GPT-4o)          3,471 ms    13%", 1),
        ("Text to speech (ElevenLabs)  19,618 ms    71%", 1),
        ("End to end                   27,457 ms", 1),
        "",
        "NFR1 required 5 seconds. Not met - and the instrumentation is what found it.",
    ], size=14, image="print/fig-29-latency-breakdown.png", image_width=Inches(3.60))
    set_notes(s[7], """
[12:00-13:15]  The most important slide after the demo. Do not soften it.

Say plainly: "The requirement was five seconds. We measure 27.5. It is not met."
Then the useful part: 71% is speech synthesis, so this is a TTS problem, not a
model problem. And the probable cause is our own prompt - it specifies 2-4
sentences and the model returns three paragraphs, so we are asking the synthesiser
to render five times more text than designed.

If asked about the "3-6 seconds" in the Phase 3 document: that number came from
the admin console's DEMO FIXTURES, not a measurement. We found that while writing
this report and corrected it. Own it - the correction is the result.

Caveat honestly: sample is 2 turns. Enough for an order of magnitude, not a
distribution.
""")

    # ---- Slide 9: challenges & limitations ------------------------------------
    set_body(s[8], [
        "Latency fails NFR1 by ~5x - dominated by speech synthesis",
        "Authorization rests on application code - RLS is enabled but unpolicied and bypassed,",
        ("and no automated test asserts isolation between users", 1),
        "Secrets committed - a long-lived token, PostHog key and Sentry DSN must be rotated",
        "No rate limiting; CORS is unrestricted",
        "Transcripts are personal data with no retention or deletion policy",
        "Coverage stops at the service layer - no widget, integration or end-to-end tests",
        "Three bundled hymns are CC BY-SA and need a credits screen the app does not have",
    ], size=15)
    set_notes(s[8], """
[13:15-14:00]  Deliver these flatly and quickly - no hedging, no "but". A limitations
slide read confidently reads as competence; one read apologetically reads as doubt.

Do not soften the secrets line. If asked "why is a token in source control?" the
honest answer is that it was a demo convenience that should have been removed, it
is in the report's pre-submission checklist, and it will be rotated.
""")

    # ---- Slide 10: conclusion & future work -----------------------------------
    set_body(s[9], [
        "Conclusion:",
        ("All 7 primary objectives met; 20/20 functional requirements implemented", 1),
        ("141 automated tests passing across both tiers, gated in CI", 1),
        ("Bilingual voice conversation working in Telugu with culturally correct address", 1),
        "Future Work:",
        ("1. Tighten the persona prompt and re-measure - likely the bulk of the 19.6 s", 1),
        ("2. Stream synthesised audio so playback starts before synthesis finishes", 1),
        ("3. Add authorization-isolation tests - the most valuable missing test", 1),
        ("4. Rotate the committed secrets; add rate limiting; restrict CORS", 1),
        ("5. Data retention policy; about/credits screen for CC BY-SA compliance", 1),
        "",
        "The system measures itself - and that is how it caught its own claim.",
    ], size=14)
    set_notes(s[9], """
[14:00-14:45]  Conclusion in three lines, future work in five, then the closing line.

The closer is the thesis of the whole project: "Most projects can tell you their app
works. This one can tell you where its seconds go - and that is how it caught a claim
it had been repeating about itself in three documents."

>>> STOP AT 14:45. Leave the closing line on screen for questions.
""")

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"  {OUT.relative_to(ROOT)}  ({OUT.stat().st_size // 1024} KB, {len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
